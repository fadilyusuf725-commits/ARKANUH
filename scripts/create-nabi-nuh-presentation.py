#!/usr/bin/env python3
"""Create CERITA NABI NUH presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os, sys

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

slides_content = [
    {"title": "Nabi Nuh dan Kaumnya", "content": "Pada zaman dahulu, hiduplah seorang nabi yang bernama Nuh. Beliau tinggal di negeri Babilonia. Nabi Nuh adalah hamba Allah yang saleh dan baik hati.\n\nKaum Nabi Nuh tidak menyembah Allah. Mereka menyembah patung-patung.", "model": "ac1160db-5473-4df9-9d05-98bce22955e3"},
    {"title": "Dakwah Nabi Nuh", "content": "Nabi Nuh sedih melihat kaumnya yang sesat. Beliau ingin mengajak mereka kembali menyembah Allah. Dengan sabar, Nabi Nuh mulai berdakwah.\n\nNabi Nuh berkata, 'Wahai kaumku, sembahlah Allah. Tidak ada Tuhan selain Dia.'", "model": "13a00c36-367b-4371-8740-b1d8bc5390b2"},
    {"title": "Penolakan Kaum", "content": "Kaum Nabi Nuh menolak ajakan beliau. Mereka berkata, 'Wahai Nuh, kamu hanya manusia biasa seperti kami!'\n\nMereka menertawakan Nabi Nuh. Nabi Nuh berdakwah selama 950 tahun lamanya!", "model": "72294589-9b8c-4554-9bf7-c6ea880c9360"},
    {"title": "Sikap Sombong", "content": "Kaum Nabi Nuh semakin sombong. Mereka berkata, 'Jika kamu benar, datangkanlah azab!'\n\nMereka menganggap pengikut Nabi Nuh hanyalah orang-orang miskin dan hina.", "model": "716b3922-1112-4456-941d-dad7f25b495e"},
    {"title": "Pengikut Setia", "content": "Hanya sedikit orang yang mau beriman kepada Nabi Nuh. Mereka adalah orang-orang yang rendah hati.\n\nMereka percaya bahwa Nabi Nuh adalah utusan Allah. Nabi Nuh sangat bersyukur memiliki pengikut yang setia.", "model": "081f023b-8053-4c46-b240-db0462e4d637"},
    {"title": "Perintah Membuat Kapal", "content": "Nabi Nuh berdoa kepada Allah, 'Ya Tuhanku, jangan biarkan seorang pun kafir tinggal di bumi.'\n\nAllah mengabulkan doa Nabi Nuh. Allah memerintahkan Nabi Nuh membuat kapal yang besar. Padahal tempat tinggalnya jauh dari laut!", "model": "228e64d9-651b-453d-b626-4e3864e069b1"},
    {"title": "Kapal Selesai Dibangun", "content": "Nabi Nuh dan pengikutnya membuat kapal besar dengan kerja keras.\n\nKaum kafir tertawa mengejek, 'Hai Nuh! Mengapa membuat kapal di padang pasir?'\n\nNabi Nuh menjawab, 'Kalian boleh mengejek sekarang. Kelak akan tahu siapa yang benar.'", "model": "2bd377ae-fbd0-4ca4-9006-ece3d2ec772c"},
    {"title": "Banjir Besar Datang", "content": "Tiba-tiba datang tanda azab dari Allah. Air memancar dari dalam tanah dengan deras.\n\nLangit gelap. Hujan belum pernah terjadi sebelumnya. Air naik ke mana-mana.\n\nNabi Nuh memerintahkan pengikutnya naik kapal bersama hewan-hewan sepasang.", "model": "6f41fe0b-7cce-40d8-8b72-93d4b01264e0"},
    {"title": "Penolakan Kan'an", "content": "Nabi Nuh memanggil putranya Kan'an. 'Wahai anakku, naiklah ke kapal!'\n\nKan'an menolak, 'Aku akan ke gunung tinggi. Gunung itu melindungiku.'\n\nNahi Nuh sedih. Anaknya tetap menolak. Sangat sedih hati beliau.", "model": "941a73af-5d63-4312-b1cb-23994ba41cc6"},
    {"title": "Semua Tenggelam", "content": "Air semakin naik. Daratan tenggelam. Orang kafir berlari ke gunung tapi air terus naik.\n\nGunung pun tenggelam. Kan'an tenggelam bersama orang kafir.\n\nKapal Nabi Nuh terapung selamat. Semua berdoa kepada Allah.", "model": "b5d24a9f-55a9-4931-b774-fead65d02b10"},
    {"title": "Kapal Berlabuh", "content": "Allah berfirman, 'Hai bumi, telanlah airmu! Hai langit, henti hujan!' Air surut.\n\nKapal Nabi Nuh berlabuh di Gunung Judi. Semua turun dengan selamat.\n\nHanya orang beriman selamat. Orang kafir sombong telah binasa.", "model": "Back cover"}
]

def add_cover_slide(prs, image_path, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    else:
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(135, 206, 250)
        bg.line.color.rgb = RGBColor(0, 0, 0)
    
    overlay = slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(6), Inches(2.625))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.fill.transparency = 0.4
    overlay.line.color.rgb = RGBColor(200, 100, 0)
    overlay.line.width = Pt(3)
    
    title_box = slide.shapes.add_textbox(Inches(2.2), Inches(1.8), Inches(5.6), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(2.2), Inches(3.1), Inches(5.6), Inches(0.8))
    stf = sub_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(28)
    sp.font.color.rgb = RGBColor(255, 255, 100)
    sp.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content, model):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
    bg.line.color.rgb = RGBColor(100, 150, 200)
    
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    ttf = tb.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(25, 25, 112)
    tp.alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(0.95), Inches(9.5), Inches(0.95))
    line.line.color.rgb = RGBColor(70, 130, 180)
    line.line.width = Pt(2)
    
    cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(3.8))
    ctf = cb.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = content
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(47, 79, 79)
    cp.space_before = Pt(6)
    cp.space_after = Pt(6)
    cp.line_spacing = 1.3
    
    mb = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9.4), Inches(0.3))
    mtf = mb.text_frame
    mtf.word_wrap = True
    mp = mtf.paragraphs[0]
    mp.text = f"🎨 3D Model: {model}"
    mp.font.size = Pt(10)
    mp.font.italic = True
    mp.font.color.rgb = RGBColor(119, 136, 153)
    
    illus = slide.shapes.add_shape(1, Inches(6.7), Inches(1.2), Inches(3), Inches(4))
    illus.fill.solid()
    illus.fill.fore_color.rgb = RGBColor(255, 255, 255)
    illus.line.color.rgb = RGBColor(176, 196, 222)
    illus.line.width = Pt(2)
    
    phb = slide.shapes.add_textbox(Inches(6.8), Inches(2.5), Inches(2.8), Inches(2))
    phtf = phb.text_frame
    phtf.word_wrap = True
    php = phtf.paragraphs[0]
    php.text = "[Illustration]\n\nGambar akan\nditambahkan"
    php.font.size = Pt(11)
    php.font.color.rgb = RGBColor(176, 176, 176)
    php.alignment = PP_ALIGN.CENTER

def create_presentation(image_path=None):
    print("📖 Creating CERITA NABI NUH presentation...\n")
    print("  ➕ Adding cover slide...")
    add_cover_slide(prs, image_path, "KISAH NABI NUH", "Media Pembelajaran")
    
    for i, data in enumerate(slides_content[:-1], 1):
        print(f"  ➕ Slide {i}/{len(slides_content)-1}: {data['title']}...")
        add_content_slide(prs, data['title'], data['content'], data['model'])
    
    print(f"  ➕ Adding back cover...")
    add_cover_slide(prs, image_path, "Terima Kasih", "Selamat Belajar!")
    
    output = r"d:\ARKANUH\public\presentations\cerita-nabi-nuh.pptx"
    os.makedirs(os.path.dirname(output), exist_ok=True)
    prs.save(output)
    
    print(f"\n✅ Created! Slides: {len(prs.slides)}")
    print(f"📍 {output}\n")

if __name__ == "__main__":
    img = sys.argv[1] if len(sys.argv) > 1 else None
    if img:
        print(f"📎 Using cover: {img}\n")
    create_presentation(img)
