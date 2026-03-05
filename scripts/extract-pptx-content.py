#!/usr/bin/env python3
"""Extract content from PPTX file and generate TypeScript data."""

import sys
import json
from pathlib import Path

# Try to import python-pptx, if not available, use zipfile fallback
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("python-pptx not available, using fallback method...", file=sys.stderr)

import zipfile
import xml.etree.ElementTree as ET


def extract_text_with_pptx(pptx_path):
    """Extract text from PPTX using python-pptx library."""
    prs = Presentation(pptx_path)
    slides_content = []
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        slides_content.append({
            "slide": slide_idx,
            "content": "\n".join(slide_text)
        })
    
    return slides_content


def extract_text_with_zipfile(pptx_path):
    """Extract text from PPTX using zipfile (fallback for when python-pptx not available)."""
    slides_content = []
    
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            # List all slide XML files
            slide_files = sorted([f for f in zip_ref.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])
            
            for slide_idx, slide_file in enumerate(slide_files, 1):
                xml_content = zip_ref.read(slide_file)
                root = ET.fromstring(xml_content)
                
                # Extract all text elements
                text_elements = []
                # Namespace for Office Open XML
                ns = {
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
                }
                
                # Find all text runs
                for t in root.findall('.//a:t', ns):
                    if t.text and t.text.strip():
                        text_elements.append(t.text.strip())
                
                content = "\n".join(text_elements)
                if content:
                    slides_content.append({
                        "slide": slide_idx,
                        "content": content
                    })
    except Exception as e:
        print(f"Error extracting with zipfile: {e}", file=sys.stderr)
        return []
    
    return slides_content


def main():
    pptx_path = Path("d:/ARKANUH/assets/flipbook/source/BUKU POP UP.pptx")
    
    if not pptx_path.exists():
        print(f"Error: PPTX file not found at {pptx_path}", file=sys.stderr)
        sys.exit(1)
    
    # Try python-pptx first, fallback to zipfile
    try:
        if HAS_PPTX:
            slides_content = extract_text_with_pptx(str(pptx_path))
        else:
            slides_content = extract_text_with_zipfile(str(pptx_path))
    except Exception as e:
        print(f"Error extracting PPTX: {e}", file=sys.stderr)
        slides_content = extract_text_with_zipfile(str(pptx_path))
    
    # Output as JSON for easy parsing
    print(json.dumps(slides_content, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
